# src/aje_libs/bd/helpers/pptx_helper.py
from pptx import Presentation
from typing import List, Dict, Optional
from ...common.logger import custom_logger

logger = custom_logger(__name__)

class PPTXHelper:
    """Helper para procesar archivos PPTX"""
    
    def __init__(self):
        pass
    
    def extract_text(self, file_path: str) -> str:
        """
        Extrae texto de una presentación PPTX.
        
        :param file_path: Ruta al archivo PPTX
        :return: Texto extraído
        """
        try:
            prs = Presentation(file_path)
            all_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text.append(shape.text)
                
                all_text.append(f"Slide {i+1}: {' '.join(slide_text)}")
            
            extracted_text = "\n\n".join(all_text)
            logger.info(f"Texto extraído exitosamente de PPTX: {file_path}")
            return extracted_text
        except Exception as e:
            logger.error(f"Error extrayendo texto de PPTX: {e}")
            raise e
    
    def extract_slides(self, file_path: str) -> List[Dict[str, str]]:
        """
        Extrae texto por diapositiva.
        
        :param file_path: Ruta al archivo PPTX
        :return: Lista de diapositivas con su contenido
        """
        try:
            prs = Presentation(file_path)
            slides = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = {
                    'slide_number': i + 1,
                    'title': '',
                    'content': []
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        # El primer shape con texto suele ser el título
                        if not slide_content['title'] and shape.text:
                            slide_content['title'] = shape.text
                        else:
                            slide_content['content'].append(shape.text)
                
                slides.append(slide_content)
            
            logger.info(f"Extraídas {len(slides)} diapositivas")
            return slides
        except Exception as e:
            logger.error(f"Error extrayendo diapositivas: {e}")
            raise e
    
    def get_slide_layouts(self, file_path: str) -> List[Dict[str, any]]:
        """
        Obtiene información sobre los layouts de las diapositivas.
        
        :param file_path: Ruta al archivo PPTX
        :return: Lista de layouts utilizados
        """
        try:
            prs = Presentation(file_path)
            layouts = []
            
            for i, slide in enumerate(prs.slides):
                layout_info = {
                    'slide_number': i + 1,
                    'layout_name': slide.slide_layout.name,
                    'shapes_count': len(slide.shapes)
                }
                layouts.append(layout_info)
            
            logger.info(f"Información de layouts obtenida para {len(layouts)} diapositivas")
            return layouts
        except Exception as e:
            logger.error(f"Error obteniendo layouts: {e}")
            raise e